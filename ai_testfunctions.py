
def process_pptx_files(filepath_list):

    from pptx import Presentation

    file_dict = {}

    for file in filepath_list:
        f = open(f'{file}', "rb")
        prs = Presentation(f)

        text_runs = []

        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text_runs.append(run.text + "\n")
    
        file_dict = " ".join(text_runs)

    return file_dict

def main():
    print("Test run of AI Test Functions Module")

if __name__ == "__main__":
    main()